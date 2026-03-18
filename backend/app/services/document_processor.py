"""
Document processor — extracts text from PDF, Excel, PPTX, and DOCX files.

Supports:
  .pdf   → pdfplumber (text + table extraction)
  .xlsx  → openpyxl  (cell values across all sheets)
  .pptx  → python-pptx (slide text and notes)
  .docx  → python-pptx paragraph reader (fallback plain text)

Returns (extracted_text: str, page_count: int | None).
"""

from __future__ import annotations

import io
import logging
from typing import Any

logger = logging.getLogger(__name__)

SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "text/plain": "txt",
}

EXTENSION_MAP = {
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".docx": "docx",
    ".doc": "docx",
    ".txt": "txt",
    ".csv": "csv",
}


def detect_file_type(filename: str, content_type: str = "") -> str | None:
    """Return a normalised file type string or None if unsupported."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in EXTENSION_MAP:
        return EXTENSION_MAP[ext]
    return SUPPORTED_TYPES.get(content_type)


def _extract_pdf(data: bytes) -> tuple[str, int]:
    import pdfplumber

    texts: list[str] = []
    page_count = 0

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            # Extract plain text
            text = page.extract_text() or ""

            # Extract tables and format as TSV-style text
            tables = page.extract_tables() or []
            table_texts: list[str] = []
            for table in tables:
                for row in table:
                    row_str = "\t".join(str(cell or "").strip() for cell in row)
                    if row_str.strip():
                        table_texts.append(row_str)

            combined = text
            if table_texts:
                combined += "\n[TABLE]\n" + "\n".join(table_texts) + "\n[/TABLE]"
            texts.append(combined)

    return "\n\n--- Page Break ---\n\n".join(texts), page_count


def _extract_xlsx(data: bytes) -> tuple[str, int]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheet_texts: list[str] = []
    sheet_count = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = [f"=== Sheet: {sheet_name} ==="]
        for row in ws.iter_rows(values_only=True):
            row_str = "\t".join(
                str(cell).strip() if cell is not None else "" for cell in row
            )
            if row_str.strip():
                rows.append(row_str)
        sheet_texts.append("\n".join(rows))

    wb.close()
    return "\n\n".join(sheet_texts), sheet_count


def _extract_pptx(data: bytes) -> tuple[str, int]:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(io.BytesIO(data))
    slide_texts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"=== Slide {i} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame else ""
            if notes:
                parts.append(f"[Notes] {notes}")
        slide_texts.append("\n".join(parts))

    return "\n\n".join(slide_texts), len(prs.slides)


def _extract_docx(data: bytes) -> tuple[str, int]:
    from docx import Document as DocxDocument  # type: ignore

    doc = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs), len(doc.paragraphs)


def _extract_txt(data: bytes) -> tuple[str, int]:
    text = data.decode("utf-8", errors="replace")
    lines = text.splitlines()
    return text, len(lines)


def _extract_csv(data: bytes) -> tuple[str, int]:
    import csv

    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(row) for row in reader]
    return "\n".join(rows), len(rows)


def extract_text(data: bytes, file_type: str) -> tuple[str, int | None]:
    """
    Extract text from uploaded file bytes.

    Returns (extracted_text, page_count).
    page_count is None for types where it's not meaningful.

    Raises ValueError if the file type is unsupported.
    Raises RuntimeError (wrapping the original error) on extraction failure.
    """
    extractors = {
        "pdf": _extract_pdf,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
        "docx": _extract_docx,
        "txt": _extract_txt,
        "csv": _extract_csv,
    }

    extractor = extractors.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type!r}")

    try:
        text, pages = extractor(data)
        # Truncate very large docs to 500k chars to keep DB sane
        if len(text) > 500_000:
            text = text[:500_000] + "\n\n[TRUNCATED — document exceeds 500k characters]"
        return text, pages
    except Exception as exc:
        logger.warning("Document extraction failed for type=%s: %s", file_type, exc)
        raise RuntimeError(f"Extraction failed: {exc}") from exc
